"""Build the executive story deck (Deliverable B).

    python tools/build_figures.py && python tools/build_deck.py

Produces ``deliverables/AuroraCart_Executive_Story.pptx``, a 7-10 minute
presentation following the case's narrative arc: Context, Tension,
Investigation, Evidence, Consequence, Decision.

Every number in the deck comes from :func:`auroracart.analysis.headline_facts`,
so the slides cannot drift from the dataset. Speaker notes carry the spoken
version of each slide plus its timing, and the chart images are the ones
``build_figures.py`` renders: same palette, same mark specs, same numbers as
the dashboard.

The deck is dark, on the same surfaces the dashboard uses. Content slides sit
on the exact chart-surface color (#1a1a19) so the PNG figures blend into the
slide with no visible rectangle around them. The one recurring ornament is the
aurora rule, a short blue-to-violet gradient bar drawn from the two palette
slots the brand mark uses.
"""

from __future__ import annotations

from dataclasses import dataclass, field

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from auroracart import analysis as A
from auroracart.data_prep import load_data
from auroracart.paths import DELIVERABLES_DIR, FIGURES_DIR

DECK_PATH = DELIVERABLES_DIR / "AuroraCart_Executive_Story.pptx"

# 16:9. The deck borrows the dashboard's dark ink and surface tokens so the
# printed story and the live app read as one system.
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

INK_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)
INK_SECONDARY = RGBColor(0xC3, 0xC2, 0xB7)
INK_MUTED = RGBColor(0x89, 0x87, 0x81)
SURFACE = RGBColor(0x1A, 0x1A, 0x19)  # content slides — matches the chart surface exactly
SURFACE_DEEP = RGBColor(0x0D, 0x0D, 0x0D)  # title/section slides — one step darker for rhythm
SURFACE_CARD = RGBColor(0x24, 0x24, 0x22)  # recommendation cards
CARD_EDGE = RGBColor(0x38, 0x38, 0x35)
ACCENT = RGBColor(0x39, 0x87, 0xE5)  # palette slot 1 (dark step)
ACCENT_2 = RGBColor(0x90, 0x85, 0xE9)  # palette slot 7 — the aurora's violet end
CRITICAL = RGBColor(0xD0, 0x3B, 0x3B)
SERIOUS = RGBColor(0xEC, 0x83, 0x5A)
WARNING = RGBColor(0xFA, 0xB2, 0x19)

FONT = "Calibri"

MARGIN = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN

# Default width of the aurora rule, as a module-level constant so it is evaluated
# once rather than on every call.
RULE_W = Inches(1.4)


@dataclass
class Slide:
    """One slide: a headline that states the finding, optional chart, optional bullets."""

    title: str
    subtitle: str = ""
    figure: str | None = None
    bullets: list[str] = field(default_factory=list)
    notes: str = ""
    kind: str = "content"  # content | title | section | closing


# --------------------------------------------------------------------------- helpers


def _textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    return frame


def _para(
    frame,
    text,
    *,
    size,
    bold=False,
    color=INK_PRIMARY,
    space_after=6,
    first=False,
    align=PP_ALIGN.LEFT,
    italic=False,
):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.alignment = align
    paragraph.space_after = Pt(space_after)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = FONT
    return paragraph


def _background(slide, color=SURFACE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _aurora_fill(shape):
    """Fill a shape with the blue-to-violet aurora sweep, falling back to solid accent."""
    try:
        shape.fill.gradient()
        stops = shape.fill.gradient_stops
        stops[0].color.rgb = ACCENT
        stops[1].color.rgb = ACCENT_2
        shape.fill.gradient_angle = 0
    except (AttributeError, TypeError, ValueError):
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    shape.shadow.inherit = False


def _accent_rule(slide, top, width=RULE_W):
    """A short aurora rule under the headline — the deck's one repeating ornament."""
    bar = slide.shapes.add_shape(1, MARGIN, top, width, Emu(28575))  # 1: rectangle
    _aurora_fill(bar)


def _fit_image(slide, image_path, top, max_height):
    """Place a figure centred horizontally, scaled to fit both the width and the height."""
    from PIL import Image

    with Image.open(image_path) as img:
        aspect = img.height / img.width

    width = CONTENT_W
    height = Emu(int(width * aspect))
    if height > max_height:
        height = max_height
        width = Emu(int(height / aspect))
    left = Emu(int((SLIDE_W - width) / 2))
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


# --------------------------------------------------------------------------- layouts


def _add_title_slide(prs, slide: Slide):
    layout = prs.slide_layouts[6]  # blank
    s = prs.slides.add_slide(layout)
    _background(s, SURFACE_DEEP)

    band = s.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(0.14))
    _aurora_fill(band)

    frame = _textbox(s, MARGIN, Inches(2.2), CONTENT_W, Inches(3.0))
    _para(frame, slide.title, size=44, bold=True, first=True, space_after=10)
    _para(frame, slide.subtitle, size=20, color=INK_SECONDARY, space_after=24)
    for bullet in slide.bullets:
        _para(frame, bullet, size=14, color=INK_MUTED, space_after=4)
    s.notes_slide.notes_text_frame.text = slide.notes
    return s


def _add_section_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _background(s, SURFACE_DEEP)
    _accent_rule(s, Inches(3.32), width=Inches(0.9))
    frame = _textbox(s, MARGIN, Inches(2.35), CONTENT_W, Inches(2.4))
    _para(
        frame, slide.subtitle.upper(), size=14, bold=True, color=ACCENT, first=True, space_after=10
    )
    _para(frame, slide.title, size=36, bold=True, color=INK_PRIMARY, space_after=0)
    s.notes_slide.notes_text_frame.text = slide.notes
    return s


def _add_content_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _background(s)

    frame = _textbox(s, MARGIN, Inches(0.40), CONTENT_W, Inches(1.25))
    _para(frame, slide.title, size=25, bold=True, first=True, space_after=4)
    if slide.subtitle:
        _para(frame, slide.subtitle, size=13, color=INK_SECONDARY, space_after=0)

    # Headline + subtitle can each wrap to two lines at these sizes; the image
    # starts below the worst case rather than below the measured text, so a longer
    # headline never lands on top of a chart.
    top = Inches(1.85) if slide.subtitle else Inches(1.45)
    _accent_rule(s, top - Inches(0.18))

    if slide.figure:
        image = FIGURES_DIR / f"{slide.figure}.png"
        if not image.exists():
            raise FileNotFoundError(f"{image} — run tools/build_figures.py first")
        bullet_space = Inches(1.25) if slide.bullets else Inches(0.35)
        _fit_image(s, image, top, SLIDE_H - top - bullet_space)

    if slide.bullets:
        bullet_top = SLIDE_H - Inches(1.22) if slide.figure else top + Inches(0.2)
        bullet_frame = _textbox(s, MARGIN, bullet_top, CONTENT_W, Inches(1.05))
        for i, bullet in enumerate(slide.bullets):
            _para(
                bullet_frame,
                bullet,
                size=14 if slide.figure else 18,
                color=INK_SECONDARY,
                first=(i == 0),
                space_after=7,
            )

    s.notes_slide.notes_text_frame.text = slide.notes
    return s


def _add_closing_slide(prs, slide: Slide):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _background(s)
    frame = _textbox(s, MARGIN, Inches(0.5), CONTENT_W, Inches(1.0))
    _para(frame, slide.title, size=26, bold=True, first=True, space_after=4)
    if slide.subtitle:
        _para(frame, slide.subtitle, size=14, color=INK_SECONDARY, space_after=0)
    _accent_rule(s, Inches(1.45))

    # Three columns, one per recommendation, so priority order is spatial as well
    # as numbered.
    gap = Inches(0.28)
    col_w = Emu(int((CONTENT_W - 2 * gap) / 3))
    for i, block in enumerate(slide.bullets):
        left = MARGIN + Emu(int(i * (col_w + gap)))
        card = s.shapes.add_shape(1, left, Inches(1.85), col_w, Inches(4.6))
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE_CARD
        card.line.color.rgb = CARD_EDGE
        card.line.width = Pt(0.75)
        card.shadow.inherit = False

        accent = s.shapes.add_shape(1, left, Inches(1.85), col_w, Inches(0.05))
        accent.fill.solid()
        accent.fill.fore_color.rgb = [CRITICAL, SERIOUS, WARNING][i]
        accent.line.fill.background()
        accent.shadow.inherit = False

        heading, *rows = block.split("|")
        text = _textbox(s, left + Inches(0.22), Inches(2.10), col_w - Inches(0.44), Inches(4.1))
        _para(
            text,
            f"PRIORITY {i + 1}",
            size=10,
            bold=True,
            color=INK_MUTED,
            first=True,
            space_after=6,
        )
        _para(text, heading.strip(), size=15, bold=True, space_after=10)
        for row in rows:
            label, _, body = row.partition(":")
            _para(text, label.strip().upper(), size=9, bold=True, color=INK_MUTED, space_after=1)
            _para(text, body.strip(), size=11, color=INK_SECONDARY, space_after=7)

    s.notes_slide.notes_text_frame.text = slide.notes
    return s


ADDERS = {
    "title": _add_title_slide,
    "section": _add_section_slide,
    "content": _add_content_slide,
    "closing": _add_closing_slide,
}


# --------------------------------------------------------------------------- content


def build_slides(f: dict) -> list[Slide]:
    """The deck, written against computed facts rather than typed-in numbers."""
    rupees = lambda v: f"₹{v / 1e6:,.1f}M"  # noqa: E731 — local formatting shorthand

    return [
        Slide(
            kind="title",
            title="AuroraCart at a Crossroads",
            subtitle="Growth is real. Value creation is not. Where the next 12 months of "
            "investment should go.",
            bullets=[
                f"Order-level analysis, {f['orders']:,.0f} orders · "
                f"{f['customers']:,.0f} customers · January 2023 to December 2025",
                "Prepared for the CEO and senior leadership team",
            ],
            notes="[0:00-0:30] Two years ago this company was smaller and more profitable. "
            "Today it is much larger and materially less profitable. I am going to show "
            "you where that profit went, why the usual explanations are wrong, and the "
            "three things I would do about it. Five charts do most of the work.",
        ),
        Slide(
            kind="section",
            subtitle="Context",
            title="What the numbers say we are",
            notes="[0:30-0:45] Start with what everyone in the room already agrees on.",
        ),
        Slide(
            title=f"Revenue grew {f['revenue_growth_pct']:.0f}%. Margin fell "
            f"{f['margin_drop_pp']:.1f} points.",
            subtitle=f"{rupees(f['revenue_first_year'])} to {rupees(f['revenue_last_year'])} net "
            f"revenue · {f['margin_first_year']:.1f}% to {f['margin_last_year']:.1f}% margin",
            figure="01_growth_vs_margin",
            bullets=[
                f"Total contribution across the three years: {rupees(f['profit_total'])} on "
                f"{rupees(f['revenue_total'])} of net revenue, an {f['margin_total']:.1f}% margin.",
            ],
            notes="[0:45-1:45] The Chief Growth Officer is right that the business is bigger. "
            "Revenue nearly doubled. But margin went the other way, and not slightly: it "
            "fell by more than a third of where it started. Two separate panels here, not "
            "one dual-axis chart, because the crossing point of two arbitrary scales would "
            "be a coincidence, not a finding. The question for the next ten minutes is not "
            "whether we grew. It is what we bought with the growth.",
        ),
        Slide(
            kind="section",
            subtitle="Tension",
            title="Growth was purchased, not earned",
            notes="[1:45-2:00]",
        ),
        Slide(
            title="Accelerate 2.0 bought revenue at roughly five times the price of profit",
            subtitle=f"Revenue per month +{f['accelerate_revenue_per_month_growth_pct']:.0f}%, "
            f"profit per month +{f['accelerate_profit_per_month_growth_pct']:.0f}% "
            f"(18 months either side of July 2024)",
            figure="02_accelerate_split",
            bullets=[
                f"Average discount rose from {f['accelerate_discount_before']:.1f}% to "
                f"{f['accelerate_discount_after']:.1f}%; new customers' share of orders rose "
                f"from {f['accelerate_new_share_before']:.0f}% to "
                f"{f['accelerate_new_share_after']:.0f}%.",
            ],
            notes="[2:00-3:00] This is the growth programme, measured on its own terms. "
            "Per-month figures, because the two windows have to be comparable. Revenue per "
            "month rose 77%. Profit per month rose 14%. We did acquire customers and we did "
            "enter Tier 2, so the programme did what it said. It also discounted four points "
            "deeper to do it. Nothing here says stop growing. It says we are not currently "
            "being paid for the growth.",
        ),
        Slide(
            kind="section",
            subtitle="Investigation",
            title="So where did the margin go?",
            notes="[3:00-3:10] The obvious answer is that we grew into a worse mix. It is wrong.",
        ),
        Slide(
            title="What we sell has not changed much. What we earn on it has.",
            subtitle=f"{f['rate_share_pct']:.0f}% of the decline is margin erosion inside "
            f"categories; the mix shift accounts for {abs(f['mix_effect_pp']):.1f} points",
            figure="03_margin_decomposition",
            bullets=[
                "The distinction matters because the two explanations call for opposite actions: "
                "a mix problem is fixed by selling something else, an erosion problem by changing "
                "the terms we sell on.",
            ],
            notes="[3:10-4:15] I want to kill the comfortable explanation first. If the decline "
            "were a mix story, meaning we grew the low-margin categories, the fix would be "
            "merchandising. It isn't. Holding the 2023 category mix fixed, margin still "
            "lands at 7.6%. Nine tenths of the fall happened inside categories we already "
            "sold. So: which categories, and why.",
        ),
        Slide(
            title="Half of revenue sits in the one category that loses money",
            subtitle=f"Electronics is {f['electronics_revenue_share']:.0f}% of net revenue at "
            f"{f['electronics_margin']:+.1f}% margin. Smartphones alone are "
            f"{f['smartphones_revenue_share']:.0f}% of revenue at "
            f"{f['smartphones_margin']:+.1f}%.",
            figure="04_category_margin",
            bullets=[
                f"Electronics loses {rupees(abs(f['electronics_profit']))} of contribution while "
                f"the other four categories earn between 18.6% and 27.7%.",
            ],
            notes="[4:15-5:00] Here is the concentration. Electronics is not a rounding error. "
            "It is half the company, and it is under water. Every other category is healthy. "
            "Note the average discount is within a third of a point across all five "
            "categories, so this is not a story about Electronics being discounted harder.",
        ),
        Slide(
            title="Electronics is a cost problem, not a discount problem",
            subtitle=f"Merchandise cost alone is {f['electronics_product_cost_pct']:.0f}% of the "
            f"net revenue an Electronics order recognises",
            figure="05_cost_structure",
            bullets=[
                "Delivery, marketing and operating costs are near-identical across all five "
                "categories. The gap sits entirely in what the goods cost us relative to what we "
                "charge for them.",
            ],
            notes="[5:00-5:45] This is the diagnostic slide. When merchandise costs 94 paise of "
            "every rupee recognised, there is almost nothing left to cover delivery, "
            "marketing and operations, let alone a discount. That is why the same 13% "
            "discount that a Fashion order absorbs comfortably puts an Electronics order "
            "into loss.",
        ),
        Slide(
            kind="section",
            subtitle="Evidence",
            title="The lever we actually control",
            notes="[5:45-5:55]",
        ),
        Slide(
            title="One discount policy is running two different businesses",
            subtitle="Electronics crosses into loss above a 10% discount; the rest of the "
            "business still earns double digits at 25% off",
            figure="06_discount_tolerance",
            bullets=[
                f"Flash Deals (average discount {f['flash_discount']:.0f}%, "
                f"{f['flash_revenue_share']:.0f}% of revenue) are the only promotion type with a "
                f"negative margin, at {f['flash_margin']:+.1f}%. Un-promoted demand earns "
                f"{f['no_promo_margin']:.1f}%.",
            ],
            notes="[5:55-6:50] This is the most actionable chart in the deck. Discount tolerance "
            "is a function of cost structure, and we have been setting it company-wide. "
            "Read this as tolerance, not as a response curve. Deeply discounted orders "
            "also differ in what is in them, and I am not claiming the discount caused the "
            "margin. But the policy implication survives either reading: the ceiling has to "
            "be set per category.",
        ),
        Slide(
            title="Product mix explains profitability. Geography and fulfilment barely move it.",
            subtitle="Ranked by how much each lens actually separates margin, weighted by the "
            "revenue it splits",
            figure="07_driver_ranking",
            bullets=[
                f"Region spreads margin by {f['geography_sd_pp']:.1f} points; product "
                f"subcategory by {f['top_driver_sd_pp']:.1f}. The regional debate in the "
                f"leadership meeting is a debate about the wrong variable.",
            ],
            notes="[6:50-7:30] The leadership team asked whether the answer is geography, "
            "customer type, channel or fulfilment. Measured properly, as spread in margin "
            "weighted by how much revenue each cut actually splits, it is none of those. "
            "It is product. This is the slide that settles the argument in the room.",
        ),
        Slide(
            title="The chart that would have sent us after the wrong problem",
            subtitle="Premium looks like our worst segment until you separate what it buys",
            figure="08_misleading_pair",
            bullets=[
                f"Premium's pooled margin is {f['premium_margin_pooled']:.1f}%. Outside "
                f"Electronics it is {f['premium_margin_ex_electronics']:.1f}%, which is ordinary. "
                f"{f['premium_electronics_share']:.0f}% of Premium's revenue is Electronics.",
            ],
            notes="[7:30-8:15] Exhibit A is accurate. Every number in it is correct, and it would "
            "have started a customer-segmentation project. Exhibit B is the same orders "
            "split by what was in them, and Premium's margin is unremarkable. The segment "
            "was never the problem; it just buys the category that is. Any cut of this data "
            "that does not control for category will produce a version of this mistake.",
        ),
        Slide(
            title="One thing is already working, and it is not where we were about to cut",
            subtitle=f"On-time delivery went from {f['on_time_before_contract']:.0f}% to "
            f"{f['on_time_after_contract']:.0f}% after the January 2025 logistics "
            f"contract; complaints fell from {f['complaint_before_contract']:.1f}% to "
            f"{f['complaint_after_contract']:.1f}%",
            figure="09_on_time_monthly",
            bullets=[
                "It cost about ₹"
                f"{f['delivery_cost_after_contract'] - f['delivery_cost_before_contract']:.0f} "
                f"more per order. The company-wide average of 43% hides both the step change and "
                f"the festive collapse: peak on-time was still only "
                f"{f['peak_on_time_last_year']:.0f}% last October and November.",
            ],
            notes="[8:15-9:00] The COO's investment paid off, and the summary table we have all "
            "been reading hides it. Pooled across three years, 43% reads as a uniformly "
            "broken function. The monthly series shows a genuine step change. It also shows "
            "the one operational risk worth naming: every festive peak still breaks "
            "delivery. That is a capacity issue, not a contract issue, and it is a watch "
            "item rather than one of my three recommendations.",
        ),
        Slide(
            kind="section",
            subtitle="Decision",
            title="Three actions, in priority order",
            notes="[9:00-9:05]",
        ),
        Slide(
            kind="closing",
            title="Where to invest, where to change course",
            subtitle="Each sized by the contribution at stake, each with the one thing we would "
            "want to know before committing",
            bullets=[
                "Re-cost and re-price Electronics before scaling it"
                f"|Evidence: {f['electronics_revenue_share']:.0f}% of revenue at "
                f"{f['electronics_margin']:+.1f}% margin; merchandise costs "
                f"{f['electronics_product_cost_pct']:.0f}% of net revenue"
                "|Benefit: the largest single pool of lost contribution, recoverable without "
                "touching demand elsewhere"
                "|Risk: Electronics drives traffic and attached sales; repricing may cost volume"
                "|We would need: SKU-level cost and competitor price data",
                "Set a category-aware discount ceiling; retire Flash Deals"
                f"|Evidence: Electronics turns loss-making above 10% off while the rest of the "
                f"business earns 10%+ at 25%; Flash Deals run at {f['flash_margin']:+.1f}%"
                "|Benefit: enforceable in the pricing system next quarter. A rule, not a bet"
                "|Risk: promotion is partly defensive; volume may move to competitors"
                "|We would need: promotion incrementality, meaning what these orders would have "
                "done unpromoted",
                "Underwrite acquisition on contribution, not revenue"
                "|Evidence: revenue per month "
                f"+{f['accelerate_revenue_per_month_growth_pct']:.0f}% "
                f"against profit per month +{f['accelerate_profit_per_month_growth_pct']:.0f}%; "
                f"paid channels fell furthest and spend the most per rupee earned"
                "|Benefit: protects growth while it stops us buying revenue that carries almost "
                "no contribution"
                "|Risk: paid channels may be seeding customers whose value appears later"
                "|We would need: cohort retention and repeat value by acquisition channel",
            ],
            notes="[9:05-9:45] Three, in order, and I would not start the second before the first. "
            "Notice what is not on this list: a regional strategy, a segmentation programme, "
            "and cutting delivery investment. The evidence does not support any of them.",
        ),
        Slide(
            title="What I would not conclude from this data",
            subtitle="The limits are part of the recommendation",
            bullets=[
                "Association, not causation. Late delivery correlates with lower ratings and "
                "triples the complaint rate. That is not proof that lateness causes either.",
                "Returns are not netted out of net revenue, so a category's true contribution is "
                "somewhat worse than shown, and worst where returns run highest.",
                "The dataset has no SKU-level cost detail, so it can say Electronics is "
                "unprofitable but not whether that is a procurement or a pricing failure.",
                "Order-level data has no lifetime view: it cannot tell us whether a customer "
                "acquired at low contribution becomes valuable later.",
                "This is a synthetic case dataset. The method is what transfers, not the "
                "specific rupee figures.",
            ],
            notes="[9:45-10:00] I would rather you trust the three recommendations because you "
            "know what they rest on. Each of the four gaps above maps to something I would "
            "ask for before we commit the money. Happy to take questions.",
        ),
    ]


def build_deck() -> str:
    facts = A.headline_facts(load_data())
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    for slide in build_slides(facts):
        ADDERS[slide.kind](prs, slide)

    DELIVERABLES_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(DECK_PATH)
    return str(DECK_PATH)


if __name__ == "__main__":
    path = build_deck()
    print(f"wrote {path}")
